#!/usr/bin/env python3
"""
Document Processing Script for n8n Integration

This script performs various document processing operations and can be executed
directly from n8n workflows using the Execute Command node.

Usage:
  python3 process_documents.py --input '{"file_path": "document.pdf", "operation": "extract_text"}'
  python3 process_documents.py --input-file input.json --operation convert
"""

import argparse
import json
import sys
import os
from pathlib import Path
from typing import Any, Dict, List, Optional, Union
import base64
import mimetypes

# Add shared libs to path
sys.path.append(str(Path(__file__).parent.parent.parent / 'libs'))
from common import (
    handle_errors, setup_logging, validate_input, safe_json_loads,
    create_success_response, create_error_response, measure_execution_time,
    safe_read_file, safe_write_file
)
from config import get_config

# Setup logging
logger = setup_logging()
config = get_config()


@measure_execution_time
@handle_errors
def process_document(
    file_path: str,
    operation: str = "extract_text",
    output_format: str = "text",
    options: Optional[Dict[str, Any]] = None
) -> Dict[str, Any]:
    """Process documents using various operations."""
    
    logger.info(f"Starting document processing: {operation} on {file_path}")
    
    # Validate file exists
    if not os.path.exists(file_path):
        return create_error_response(
            f"File not found: {file_path}",
            "FileNotFoundError",
            {"file_path": file_path}
        )
    
    # Get file info
    file_info = get_file_info(file_path)
    
    options = options or {}
    
    if operation == "extract_text":
        return extract_text_from_document(file_path, file_info, options)
    elif operation == "extract_metadata":
        return extract_document_metadata(file_path, file_info, options)
    elif operation == "convert":
        return convert_document(file_path, file_info, output_format, options)
    elif operation == "split":
        return split_document(file_path, file_info, options)
    elif operation == "merge":
        return merge_documents(file_path, file_info, options)
    elif operation == "extract_images":
        return extract_images_from_document(file_path, file_info, options)
    elif operation == "ocr":
        return perform_ocr(file_path, file_info, options)
    elif operation == "analyze":
        return analyze_document(file_path, file_info, options)
    elif operation == "compress":
        return compress_document(file_path, file_info, options)
    elif operation == "watermark":
        return add_watermark(file_path, file_info, options)
    else:
        return create_error_response(
            f"Unknown operation: {operation}",
            "ValueError",
            {"available_operations": [
                "extract_text", "extract_metadata", "convert", "split", 
                "merge", "extract_images", "ocr", "analyze", "compress", "watermark"
            ]}
        )


def get_file_info(file_path: str) -> Dict[str, Any]:
    """Get basic file information."""
    
    stat = os.stat(file_path)
    mime_type, _ = mimetypes.guess_type(file_path)
    
    return {
        "path": file_path,
        "name": os.path.basename(file_path),
        "size": stat.st_size,
        "extension": Path(file_path).suffix.lower(),
        "mime_type": mime_type,
        "modified_time": stat.st_mtime,
        "created_time": stat.st_ctime
    }


def extract_text_from_document(file_path: str, file_info: Dict[str, Any], options: Dict[str, Any]) -> Dict[str, Any]:
    """Extract text from various document formats."""
    
    extension = file_info["extension"]
    
    try:
        if extension == ".pdf":
            return extract_text_from_pdf(file_path, options)
        elif extension in [".docx", ".doc"]:
            return extract_text_from_word(file_path, options)
        elif extension in [".xlsx", ".xls"]:
            return extract_text_from_excel(file_path, options)
        elif extension in [".pptx", ".ppt"]:
            return extract_text_from_powerpoint(file_path, options)
        elif extension in [".txt", ".md", ".csv"]:
            return extract_text_from_plain_text(file_path, options)
        elif extension in [".html", ".htm"]:
            return extract_text_from_html(file_path, options)
        elif extension in [".rtf"]:
            return extract_text_from_rtf(file_path, options)
        else:
            return create_error_response(
                f"Unsupported file format for text extraction: {extension}",
                "ValueError",
                {"supported_formats": [".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt", ".txt", ".md", ".csv", ".html", ".htm", ".rtf"]}
            )
    
    except Exception as e:
        logger.error(f"Error extracting text from {file_path}: {e}")
        return create_error_response(
            f"Text extraction failed: {str(e)}",
            type(e).__name__
        )


def extract_text_from_pdf(file_path: str, options: Dict[str, Any]) -> Dict[str, Any]:
    """Extract text from PDF files."""
    
    try:
        import PyPDF2
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            text_content = []
            page_texts = []
            
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    page_texts.append({
                        "page_number": page_num + 1,
                        "text": page_text,
                        "word_count": len(page_text.split())
                    })
                    text_content.append(page_text)
                except Exception as e:
                    logger.warning(f"Error extracting text from page {page_num + 1}: {e}")
                    page_texts.append({
                        "page_number": page_num + 1,
                        "text": "",
                        "word_count": 0,
                        "error": str(e)
                    })
            
            full_text = "\n\n".join(text_content)
            
            # Extract metadata
            metadata = {}
            if pdf_reader.metadata:
                for key, value in pdf_reader.metadata.items():
                    if key.startswith('/'):
                        key = key[1:]  # Remove leading slash
                    metadata[key.lower()] = str(value) if value else None
            
            result = {
                "full_text": full_text,
                "pages": page_texts,
                "page_count": len(pdf_reader.pages),
                "word_count": len(full_text.split()),
                "character_count": len(full_text),
                "metadata": metadata
            }
            
            return create_success_response(result, {
                "operation": "extract_text",
                "format": "pdf",
                "pages_processed": len(page_texts)
            })
    
    except ImportError:
        return create_error_response(
            "PDF processing requires PyPDF2 library",
            "ImportError",
            {"required_packages": ["PyPDF2"]}
        )
    except Exception as e:
        logger.error(f"Error processing PDF: {e}")
        return create_error_response(
            f"PDF processing failed: {str(e)}",
            type(e).__name__
        )


def extract_text_from_word(file_path: str, options: Dict[str, Any]) -> Dict[str, Any]:
    """Extract text from Word documents."""
    
    try:
        from docx import Document
        
        doc = Document(file_path)
        
        # Extract paragraphs
        paragraphs = []
        full_text_parts = []
        
        for i, paragraph in enumerate(doc.paragraphs):
            text = paragraph.text.strip()
            if text:  # Skip empty paragraphs
                paragraphs.append({
                    "paragraph_number": i + 1,
                    "text": text,
                    "word_count": len(text.split())
                })
                full_text_parts.append(text)
        
        # Extract tables
        tables = []
        for table_num, table in enumerate(doc.tables):
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            
            tables.append({
                "table_number": table_num + 1,
                "rows": len(table_data),
                "columns": len(table_data[0]) if table_data else 0,
                "data": table_data
            })
        
        full_text = "\n\n".join(full_text_parts)
        
        # Extract core properties
        core_props = doc.core_properties
        metadata = {
            "title": core_props.title,
            "author": core_props.author,
            "subject": core_props.subject,
            "keywords": core_props.keywords,
            "created": core_props.created.isoformat() if core_props.created else None,
            "modified": core_props.modified.isoformat() if core_props.modified else None,
            "last_modified_by": core_props.last_modified_by,
            "revision": core_props.revision
        }
        
        result = {
            "full_text": full_text,
            "paragraphs": paragraphs,
            "tables": tables,
            "paragraph_count": len(paragraphs),
            "table_count": len(tables),
            "word_count": len(full_text.split()),
            "character_count": len(full_text),
            "metadata": metadata
        }
        
        return create_success_response(result, {
            "operation": "extract_text",
            "format": "word",
            "paragraphs_processed": len(paragraphs),
            "tables_processed": len(tables)
        })
    
    except ImportError:
        return create_error_response(
            "Word document processing requires python-docx library",
            "ImportError",
            {"required_packages": ["python-docx"]}
        )
    except Exception as e:
        logger.error(f"Error processing Word document: {e}")
        return create_error_response(
            f"Word document processing failed: {str(e)}",
            type(e).__name__
        )


def extract_text_from_excel(file_path: str, options: Dict[str, Any]) -> Dict[str, Any]:
    """Extract text from Excel files."""
    
    try:
        import pandas as pd
        
        # Read all sheets
        excel_file = pd.ExcelFile(file_path)
        sheets_data = {}
        all_text_parts = []
        
        for sheet_name in excel_file.sheet_names:
            try:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Convert to text
                sheet_text = df.to_string(index=False)
                all_text_parts.append(f"Sheet: {sheet_name}\n{sheet_text}")
                
                # Store sheet info
                sheets_data[sheet_name] = {
                    "rows": len(df),
                    "columns": len(df.columns),
                    "column_names": df.columns.tolist(),
                    "text": sheet_text,
                    "data": df.to_dict('records') if len(df) <= 100 else "Too many rows to include"
                }
            
            except Exception as e:
                logger.warning(f"Error processing sheet {sheet_name}: {e}")
                sheets_data[sheet_name] = {"error": str(e)}
        
        full_text = "\n\n".join(all_text_parts)
        
        result = {
            "full_text": full_text,
            "sheets": sheets_data,
            "sheet_count": len(excel_file.sheet_names),
            "word_count": len(full_text.split()),
            "character_count": len(full_text)
        }
        
        return create_success_response(result, {
            "operation": "extract_text",
            "format": "excel",
            "sheets_processed": len(sheets_data)
        })
    
    except ImportError:
        return create_error_response(
            "Excel processing requires pandas and openpyxl libraries",
            "ImportError",
            {"required_packages": ["pandas", "openpyxl"]}
        )
    except Exception as e:
        logger.error(f"Error processing Excel file: {e}")
        return create_error_response(
            f"Excel processing failed: {str(e)}",
            type(e).__name__
        )


def extract_text_from_powerpoint(file_path: str, options: Dict[str, Any]) -> Dict[str, Any]:
    """Extract text from PowerPoint files."""
    
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        
        slides_data = []
        all_text_parts = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text_parts = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_parts.append(shape.text.strip())
            
            slide_text = "\n".join(slide_text_parts)
            all_text_parts.append(slide_text)
            
            slides_data.append({
                "slide_number": slide_num + 1,
                "text": slide_text,
                "word_count": len(slide_text.split()),
                "shape_count": len(slide.shapes)
            })
        
        full_text = "\n\n".join(all_text_parts)
        
        result = {
            "full_text": full_text,
            "slides": slides_data,
            "slide_count": len(prs.slides),
            "word_count": len(full_text.split()),
            "character_count": len(full_text)
        }
        
        return create_success_response(result, {
            "operation": "extract_text",
            "format": "powerpoint",
            "slides_processed": len(slides_data)
        })
    
    except ImportError:
        return create_error_response(
            "PowerPoint processing requires python-pptx library",
            "ImportError",
            {"required_packages": ["python-pptx"]}
        )
    except Exception as e:
        logger.error(f"Error processing PowerPoint file: {e}")
        return create_error_response(
            f"PowerPoint processing failed: {str(e)}",
            type(e).__name__
        )


def extract_text_from_plain_text(file_path: str, options: Dict[str, Any]) -> Dict[str, Any]:
    """Extract text from plain text files."""
    
    try:
        encoding = options.get("encoding", "utf-8")
        
        with open(file_path, 'r', encoding=encoding) as file:
            content = file.read()
        
        # Split into lines
        lines = content.split('\n')
        
        result = {
            "full_text": content,
            "lines": [{
                "line_number": i + 1,
                "text": line,
                "word_count": len(line.split())
            } for i, line in enumerate(lines)],
            "line_count": len(lines),
            "word_count": len(content.split()),
            "character_count": len(content),
            "encoding": encoding
        }
        
        return create_success_response(result, {
            "operation": "extract_text",
            "format": "plain_text",
            "lines_processed": len(lines)
        })
    
    except UnicodeDecodeError as e:
        return create_error_response(
            f"Encoding error: {str(e)}. Try specifying a different encoding.",
            "UnicodeDecodeError",
            {"suggested_encodings": ["utf-8", "latin-1", "cp1252"]}
        )
    except Exception as e:
        logger.error(f"Error processing plain text file: {e}")
        return create_error_response(
            f"Plain text processing failed: {str(e)}",
            type(e).__name__
        )


def extract_text_from_html(file_path: str, options: Dict[str, Any]) -> Dict[str, Any]:
    """Extract text from HTML files."""
    
    try:
        from bs4 import BeautifulSoup
        
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        soup = BeautifulSoup(content, 'html.parser')
        
        # Extract text
        text = soup.get_text(separator='\n', strip=True)
        
        # Extract specific elements
        title = soup.title.string if soup.title else None
        headings = []
        
        for i in range(1, 7):
            for heading in soup.find_all(f'h{i}'):
                headings.append({
                    "level": i,
                    "text": heading.get_text(strip=True)
                })
        
        # Extract links
        links = []
        for link in soup.find_all('a', href=True):
            links.append({
                "text": link.get_text(strip=True),
                "href": link['href']
            })
        
        result = {
            "full_text": text,
            "title": title,
            "headings": headings,
            "links": links,
            "word_count": len(text.split()),
            "character_count": len(text),
            "heading_count": len(headings),
            "link_count": len(links)
        }
        
        return create_success_response(result, {
            "operation": "extract_text",
            "format": "html",
            "elements_extracted": {
                "headings": len(headings),
                "links": len(links)
            }
        })
    
    except ImportError:
        return create_error_response(
            "HTML processing requires beautifulsoup4 library",
            "ImportError",
            {"required_packages": ["beautifulsoup4"]}
        )
    except Exception as e:
        logger.error(f"Error processing HTML file: {e}")
        return create_error_response(
            f"HTML processing failed: {str(e)}",
            type(e).__name__
        )


def extract_text_from_rtf(file_path: str, options: Dict[str, Any]) -> Dict[str, Any]:
    """Extract text from RTF files."""
    
    try:
        from striprtf.striprtf import rtf_to_text
        
        with open(file_path, 'r', encoding='utf-8') as file:
            rtf_content = file.read()
        
        text = rtf_to_text(rtf_content)
        
        result = {
            "full_text": text,
            "word_count": len(text.split()),
            "character_count": len(text)
        }
        
        return create_success_response(result, {
            "operation": "extract_text",
            "format": "rtf"
        })
    
    except ImportError:
        return create_error_response(
            "RTF processing requires striprtf library",
            "ImportError",
            {"required_packages": ["striprtf"]}
        )
    except Exception as e:
        logger.error(f"Error processing RTF file: {e}")
        return create_error_response(
            f"RTF processing failed: {str(e)}",
            type(e).__name__
        )


def extract_document_metadata(file_path: str, file_info: Dict[str, Any], options: Dict[str, Any]) -> Dict[str, Any]:
    """Extract metadata from documents."""
    
    try:
        from exifread import process_file
        
        metadata = {
            "file_info": file_info,
            "exif_data": {},
            "document_properties": {}
        }
        
        # Try to extract EXIF data for images
        if file_info["mime_type"] and file_info["mime_type"].startswith("image/"):
            try:
                with open(file_path, 'rb') as f:
                    tags = process_file(f)
                    for tag in tags.keys():
                        metadata["exif_data"][tag] = str(tags[tag])
            except Exception as e:
                logger.warning(f"Could not extract EXIF data: {e}")
        
        # Extract document-specific metadata
        extension = file_info["extension"]
        
        if extension == ".pdf":
            try:
                import PyPDF2
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    if pdf_reader.metadata:
                        for key, value in pdf_reader.metadata.items():
                            if key.startswith('/'):
                                key = key[1:]
                            metadata["document_properties"][key.lower()] = str(value) if value else None
            except Exception as e:
                logger.warning(f"Could not extract PDF metadata: {e}")
        
        elif extension in [".docx", ".doc"]:
            try:
                from docx import Document
                doc = Document(file_path)
                core_props = doc.core_properties
                metadata["document_properties"] = {
                    "title": core_props.title,
                    "author": core_props.author,
                    "subject": core_props.subject,
                    "keywords": core_props.keywords,
                    "created": core_props.created.isoformat() if core_props.created else None,
                    "modified": core_props.modified.isoformat() if core_props.modified else None,
                    "last_modified_by": core_props.last_modified_by,
                    "revision": core_props.revision
                }
            except Exception as e:
                logger.warning(f"Could not extract Word metadata: {e}")
        
        return create_success_response(metadata, {
            "operation": "extract_metadata",
            "format": extension,
            "metadata_fields": len(metadata["document_properties"]) + len(metadata["exif_data"])
        })
    
    except ImportError:
        return create_error_response(
            "Metadata extraction requires exifread library",
            "ImportError",
            {"required_packages": ["exifread"]}
        )
    except Exception as e:
        logger.error(f"Error extracting metadata: {e}")
        return create_error_response(
            f"Metadata extraction failed: {str(e)}",
            type(e).__name__
        )


def convert_document(file_path: str, file_info: Dict[str, Any], output_format: str, options: Dict[str, Any]) -> Dict[str, Any]:
    """Convert documents between formats."""
    
    # This is a placeholder for document conversion
    # In a real implementation, you would use libraries like:
    # - pandoc for general document conversion
    # - unoconv for LibreOffice conversions
    # - specific libraries for format conversions
    
    return create_error_response(
        "Document conversion not implemented yet",
        "NotImplementedError",
        {
            "suggested_tools": ["pandoc", "unoconv", "libreoffice"],
            "input_format": file_info["extension"],
            "output_format": output_format
        }
    )


def split_document(file_path: str, file_info: Dict[str, Any], options: Dict[str, Any]) -> Dict[str, Any]:
    """Split documents into smaller parts."""
    
    # Placeholder for document splitting functionality
    return create_error_response(
        "Document splitting not implemented yet",
        "NotImplementedError"
    )


def merge_documents(file_path: str, file_info: Dict[str, Any], options: Dict[str, Any]) -> Dict[str, Any]:
    """Merge multiple documents."""
    
    # Placeholder for document merging functionality
    return create_error_response(
        "Document merging not implemented yet",
        "NotImplementedError"
    )


def extract_images_from_document(file_path: str, file_info: Dict[str, Any], options: Dict[str, Any]) -> Dict[str, Any]:
    """Extract images from documents."""
    
    # Placeholder for image extraction functionality
    return create_error_response(
        "Image extraction not implemented yet",
        "NotImplementedError"
    )


def perform_ocr(file_path: str, file_info: Dict[str, Any], options: Dict[str, Any]) -> Dict[str, Any]:
    """Perform OCR on images or scanned documents."""
    
    try:
        import pytesseract
        from PIL import Image
        
        # Open image
        image = Image.open(file_path)
        
        # Perform OCR
        language = options.get("language", "eng")
        text = pytesseract.image_to_string(image, lang=language)
        
        # Get additional OCR data
        data = pytesseract.image_to_data(image, lang=language, output_type=pytesseract.Output.DICT)
        
        # Process OCR data
        words = []
        for i in range(len(data['text'])):
            if int(data['conf'][i]) > 0:  # Only include words with confidence > 0
                words.append({
                    "text": data['text'][i],
                    "confidence": int(data['conf'][i]),
                    "left": int(data['left'][i]),
                    "top": int(data['top'][i]),
                    "width": int(data['width'][i]),
                    "height": int(data['height'][i])
                })
        
        result = {
            "extracted_text": text,
            "words": words,
            "word_count": len(text.split()),
            "character_count": len(text),
            "language": language,
            "average_confidence": sum(w["confidence"] for w in words) / len(words) if words else 0
        }
        
        return create_success_response(result, {
            "operation": "ocr",
            "words_detected": len(words),
            "language_used": language
        })
    
    except ImportError:
        return create_error_response(
            "OCR requires pytesseract and Pillow libraries",
            "ImportError",
            {"required_packages": ["pytesseract", "Pillow"]}
        )
    except Exception as e:
        logger.error(f"Error performing OCR: {e}")
        return create_error_response(
            f"OCR failed: {str(e)}",
            type(e).__name__
        )


def analyze_document(file_path: str, file_info: Dict[str, Any], options: Dict[str, Any]) -> Dict[str, Any]:
    """Analyze document structure and content."""
    
    # First extract text
    text_result = extract_text_from_document(file_path, file_info, options)
    
    if not text_result.get("success"):
        return text_result
    
    text_data = text_result["data"]
    full_text = text_data.get("full_text", "")
    
    if not full_text:
        return create_error_response(
            "No text content found for analysis",
            "ValueError"
        )
    
    # Perform text analysis
    analysis = {
        "text_statistics": {
            "word_count": len(full_text.split()),
            "character_count": len(full_text),
            "sentence_count": len([s for s in full_text.split('.') if s.strip()]),
            "paragraph_count": len([p for p in full_text.split('\n\n') if p.strip()]),
            "average_words_per_sentence": 0,
            "average_characters_per_word": 0
        },
        "readability": {},
        "language_analysis": {},
        "content_structure": text_data
    }
    
    # Calculate averages
    words = full_text.split()
    sentences = [s for s in full_text.split('.') if s.strip()]
    
    if sentences:
        analysis["text_statistics"]["average_words_per_sentence"] = len(words) / len(sentences)
    
    if words:
        analysis["text_statistics"]["average_characters_per_word"] = sum(len(word) for word in words) / len(words)
    
    # Try to add readability analysis
    try:
        import textstat
        analysis["readability"] = {
            "flesch_reading_ease": textstat.flesch_reading_ease(full_text),
            "flesch_kincaid_grade": textstat.flesch_kincaid_grade(full_text),
            "gunning_fog": textstat.gunning_fog(full_text)
        }
    except ImportError:
        logger.warning("textstat not available for readability analysis")
    
    # Try to add language detection
    try:
        from langdetect import detect
        analysis["language_analysis"]["detected_language"] = detect(full_text)
    except ImportError:
        logger.warning("langdetect not available for language analysis")
    except Exception:
        logger.warning("Could not detect language")
    
    return create_success_response(analysis, {
        "operation": "analyze",
        "format": file_info["extension"],
        "analysis_components": list(analysis.keys())
    })


def compress_document(file_path: str, file_info: Dict[str, Any], options: Dict[str, Any]) -> Dict[str, Any]:
    """Compress documents."""
    
    # Placeholder for document compression functionality
    return create_error_response(
        "Document compression not implemented yet",
        "NotImplementedError"
    )


def add_watermark(file_path: str, file_info: Dict[str, Any], options: Dict[str, Any]) -> Dict[str, Any]:
    """Add watermark to documents."""
    
    # Placeholder for watermarking functionality
    return create_error_response(
        "Document watermarking not implemented yet",
        "NotImplementedError"
    )


def main():
    """Main function for command-line usage."""
    
    parser = argparse.ArgumentParser(
        description="Process documents using various operations",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Extract text from PDF
  python3 process_documents.py --input '{"file_path": "document.pdf", "operation": "extract_text"}'
  
  # Extract metadata
  python3 process_documents.py --input '{"file_path": "document.docx", "operation": "extract_metadata"}'
  
  # Perform OCR on image
  python3 process_documents.py --input '{"file_path": "scan.png", "operation": "ocr", "options": {"language": "eng"}}'
"""
    )
    
    # Input options
    input_group = parser.add_mutually_exclusive_group(required=True)
    input_group.add_argument('--input', help='JSON input data as string')
    input_group.add_argument('--input-file', help='Path to JSON input file')
    
    # Operation options
    parser.add_argument(
        '--operation', 
        default='extract_text',
        choices=[
            'extract_text', 'extract_metadata', 'convert', 'split', 
            'merge', 'extract_images', 'ocr', 'analyze', 'compress', 'watermark'
        ],
        help='Document processing operation (default: extract_text)'
    )
    
    # Output options
    parser.add_argument('--output-file', help='Path to save output JSON file')
    parser.add_argument('--pretty', action='store_true', help='Pretty print JSON output')
    
    args = parser.parse_args()
    
    try:
        # Parse input data
        if args.input:
            input_data = safe_json_loads(args.input)
        else:
            with open(args.input_file, 'r', encoding='utf-8') as f:
                input_data = json.load(f)
        
        # Validate input structure
        schema = {
            "file_path": {"type": "string", "required": True},
            "operation": {"type": "string", "required": False},
            "output_format": {"type": "string", "required": False},
            "options": {"type": "object", "required": False}
        }
        
        validate_input(input_data, schema)
        
        # Extract parameters
        file_path = input_data["file_path"]
        operation = input_data.get("operation", args.operation)
        output_format = input_data.get("output_format", "text")
        options = input_data.get("options", {})
        
        # Process document
        result = process_document(
            file_path=file_path,
            operation=operation,
            output_format=output_format,
            options=options
        )
        
        # Output result
        output_json = json.dumps(result, indent=2 if args.pretty else None, ensure_ascii=False)
        
        if args.output_file:
            with open(args.output_file, 'w', encoding='utf-8') as f:
                f.write(output_json)
            logger.info(f"Results saved to {args.output_file}")
        else:
            print(output_json)
    
    except Exception as e:
        error_result = create_error_response(str(e), type(e).__name__)
        print(json.dumps(error_result), file=sys.stderr)
        sys.exit(1)


if __name__ == '__main__':
    main()